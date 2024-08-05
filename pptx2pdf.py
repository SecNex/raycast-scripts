#!/usr/bin/env python3

# Required parameters:
# @raycast.schemaVersion 1
# @raycast.title PPT2PDF
# @raycast.mode silent

# Optional parameters:
# @raycast.icon 📊
# @raycast.argument1 { "type": "text", "placeholder": "Path to PPT" }
# @raycast.packageName Quick Tools

# Documentation:
# @raycast.author SecNex
# @raycast.authorURL https://docs.secnex.io

import os
import sys
from pptx import Presentation
from fpdf import FPDF

print("PPT2PDF - Convert PowerPoint files to PDF files")

def convert_ppt_to_pdf(input_file, output_file):
    try:
        # Load presentation
        prs = Presentation(input_file)
        
        # Create a new PDF document
        pdf = FPDF()

        # Loop through slides
        for slide_number, slide in enumerate(prs.slides):
            # Save slide as image
            slide_path = f"slide_{slide_number}.png"
            slide.shapes._spTree.save(slide_path)
            
            # Convert slide image to PDF
            pdf.add_page()
            pdf.image(slide_path, 0, 0, 210, 297)
            os.remove(slide_path)

        # Save PDF
        pdf.output(output_file)
        print(f"Converted {input_file} to {output_file} successfully.")
    except Exception as e:
        print(f"Failed to convert {input_file}: {e}")

def process_path(path):
    if os.path.isfile(path) and path.lower().endswith(('.ppt', '.pptx')):
        print(f"Processing single file: {path}")
        output_path = os.path.splitext(path)[0] + '.pdf'
        convert_ppt_to_pdf(path, output_path)
    else:
        print(f"The path {path} is not a valid PowerPoint file.")

if __name__ == "__main__":
    path = sys.argv[1]
    try:
        process_path(path=path)
    except Exception as e:
        print(f"Failed to convert {path}: {e}")
        sys.exit(1)